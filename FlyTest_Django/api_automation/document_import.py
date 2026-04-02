from __future__ import annotations

import importlib.util
import json
import os
import random
import re
import shlex
import shutil
import subprocess
import tempfile
import uuid
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any
from urllib.parse import parse_qsl, urljoin

import yaml
from bs4 import BeautifulSoup


HTTP_METHODS = {"GET", "POST", "PUT", "PATCH", "DELETE", "HEAD", "OPTIONS"}
STRUCTURED_EXTENSIONS = {".json", ".yaml", ".yml"}
TEXT_EXTENSIONS = {".txt", ".md"}
NATIVE_DOCUMENT_EXTENSIONS = {".docx", ".pptx", ".xlsx", ".epub", ".html", ".htm"}
MARKER_EXTENSIONS = {
    ".pdf",
    ".png",
    ".jpg",
    ".jpeg",
    ".webp",
    ".bmp",
    ".tif",
    ".tiff",
    ".doc",
    ".xls",
}

PLACEHOLDER_PATTERN = re.compile(r"\{\{\s*([a-zA-Z0-9_.-]+)\s*\}\}")
ENDPOINT_PATTERN = re.compile(
    r"(?im)^(?:#{1,6}\s+)?(?P<method>GET|POST|PUT|PATCH|DELETE|HEAD|OPTIONS)\s+(?P<url>(?:https?://|/)[^\s`]+)"
)
HEADING_PATTERN = re.compile(r"(?m)^(#{1,6})\s+(?P<title>.+)$")
CODE_BLOCK_PATTERN = re.compile(r"```(?P<lang>[a-zA-Z0-9_-]*)\n(?P<content>.*?)```", re.S)
URL_IN_TEXT_PATTERN = re.compile(r"((?:https?://|/)[^\s`\"'<>]+)")
STATUS_CODE_TEXT_PATTERN = re.compile(r"(?:status(?:\s*code)?|状态码|响应码)\s*[:：]?\s*(\d{3})", re.I)
JSON_PATH_TEXT_PATTERN = re.compile(r"((?:data|result|response)(?:\.[A-Za-z0-9_-]+)+)")
DEFAULT_REQUEST_TIMEOUT_MS = 30000
JSON_CONTENT_TYPES = {"application/json"}
XML_CONTENT_TYPES = {"application/xml", "text/xml"}
TEXT_CONTENT_TYPES = {"text/plain"}


@dataclass
class ParsedRequestData:
    name: str
    method: str
    url: str
    description: str = ""
    headers: dict[str, Any] = field(default_factory=dict)
    params: dict[str, Any] = field(default_factory=dict)
    body_type: str = "none"
    body: Any = field(default_factory=dict)
    assertions: list[dict[str, Any]] = field(default_factory=list)
    request_spec: dict[str, Any] | None = None
    assertion_specs: list[dict[str, Any]] = field(default_factory=list)
    extractor_specs: list[dict[str, Any]] = field(default_factory=list)
    collection_name: str | None = None


@dataclass
class DocumentImportResult:
    source_type: str
    requests: list[ParsedRequestData]
    marker_used: bool = False
    note: str = ""


def default_auth_spec() -> dict[str, Any]:
    return {
        "auth_type": "none",
        "username": "",
        "password": "",
        "token_value": "",
        "token_variable": "token",
        "header_name": "Authorization",
        "bearer_prefix": "Bearer",
        "api_key_name": "",
        "api_key_in": "header",
        "api_key_value": "",
        "cookie_name": "",
        "bootstrap_request_id": None,
        "bootstrap_request_name": "",
        "bootstrap_token_path": "",
    }


def default_transport_spec() -> dict[str, Any]:
    return {
        "verify_ssl": True,
        "proxy_url": "",
        "client_cert": "",
        "client_key": "",
        "follow_redirects": True,
        "retry_count": 0,
        "retry_interval_ms": 500,
    }


def named_items_from_mapping(mapping: dict[str, Any] | None) -> list[dict[str, Any]]:
    if not isinstance(mapping, dict):
        return []
    return [
        {
            "name": str(name),
            "value": value,
            "enabled": True,
            "order": index,
        }
        for index, (name, value) in enumerate(mapping.items())
    ]


def build_request_spec(
    *,
    method: str,
    url: str,
    headers: dict[str, Any] | None = None,
    query: dict[str, Any] | None = None,
    cookies: dict[str, Any] | None = None,
    body_mode: str = "none",
    body_json: Any | None = None,
    raw_text: str = "",
    xml_text: str = "",
    binary_base64: str = "",
    graphql_query: str = "",
    graphql_operation_name: str = "",
    graphql_variables: dict[str, Any] | None = None,
    form_fields: dict[str, Any] | None = None,
    multipart_parts: dict[str, Any] | None = None,
    files: list[dict[str, Any]] | None = None,
    auth: dict[str, Any] | None = None,
    transport: dict[str, Any] | None = None,
    timeout_ms: int = DEFAULT_REQUEST_TIMEOUT_MS,
) -> dict[str, Any]:
    return {
        "method": method,
        "url": url,
        "body_mode": body_mode,
        "body_json": body_json if isinstance(body_json, (dict, list)) else {},
        "raw_text": raw_text,
        "xml_text": xml_text,
        "binary_base64": binary_base64,
        "graphql_query": graphql_query,
        "graphql_operation_name": graphql_operation_name,
        "graphql_variables": graphql_variables if isinstance(graphql_variables, dict) else {},
        "timeout_ms": timeout_ms,
        "headers": named_items_from_mapping(headers),
        "query": named_items_from_mapping(query),
        "cookies": named_items_from_mapping(cookies),
        "form_fields": named_items_from_mapping(form_fields),
        "multipart_parts": named_items_from_mapping(multipart_parts),
        "files": files or [],
        "auth": auth or default_auth_spec(),
        "transport": transport or default_transport_spec(),
    }


def request_spec_to_legacy_body(spec: dict[str, Any]) -> tuple[str, Any]:
    body_mode = str(spec.get("body_mode") or "none").lower()
    def to_mapping(items: Any) -> dict[str, Any]:
        if isinstance(items, dict):
            return items
        if isinstance(items, list):
            return {
                item["name"]: item.get("value", "")
                for item in items
                if isinstance(item, dict) and item.get("name")
            }
        return {}
    if body_mode == "json":
        return "json", spec.get("body_json") or {}
    if body_mode in {"form", "urlencoded"}:
        body = to_mapping(spec.get("form_fields"))
        return "form", body
    if body_mode == "multipart":
        body = to_mapping(spec.get("multipart_parts"))
        return "form", body
    if body_mode == "graphql":
        return "json", {
            "query": spec.get("graphql_query") or "",
            "operationName": spec.get("graphql_operation_name") or "",
            "variables": spec.get("graphql_variables") or {},
        }
    if body_mode == "xml":
        return "raw", spec.get("xml_text") or ""
    if body_mode == "binary":
        return "raw", spec.get("binary_base64") or ""
    if body_mode == "raw":
        return "raw", spec.get("raw_text") or ""
    return "none", {}


def is_json_content_type(content_type: str) -> bool:
    lowered = content_type.lower()
    return lowered in JSON_CONTENT_TYPES or lowered.endswith("+json")


def is_xml_content_type(content_type: str) -> bool:
    lowered = content_type.lower()
    return lowered in XML_CONTENT_TYPES or lowered.endswith("+xml")


def is_graphql_content_type(content_type: str) -> bool:
    return "graphql" in content_type.lower()


def is_binary_schema(schema: dict[str, Any] | None) -> bool:
    if not isinstance(schema, dict):
        return False
    if schema.get("format") in {"binary", "base64"} or schema.get("contentEncoding") == "base64":
        return True
    if schema.get("type") == "array":
        return is_binary_schema(schema.get("items") if isinstance(schema.get("items"), dict) else None)
    return False


def _resolve_multipart_content_type(
    field_name: str,
    schema: dict[str, Any] | None,
    encoding: dict[str, Any] | None = None,
) -> str:
    encoding_item = encoding.get(field_name) if isinstance(encoding, dict) else None
    if isinstance(encoding_item, dict) and encoding_item.get("contentType"):
        return str(encoding_item.get("contentType"))

    if not isinstance(schema, dict):
        return "application/octet-stream"
    if schema.get("contentMediaType"):
        return str(schema.get("contentMediaType"))
    if schema.get("type") == "array":
        item_schema = schema.get("items") if isinstance(schema.get("items"), dict) else {}
        if isinstance(item_schema, dict) and item_schema.get("contentMediaType"):
            return str(item_schema.get("contentMediaType"))
    return "application/octet-stream"


def _build_multipart_file_entries(
    *,
    field_name: str,
    value: Any,
    schema: dict[str, Any] | None = None,
    encoding: dict[str, Any] | None = None,
    order_start: int = 0,
) -> list[dict[str, Any]]:
    values = value if isinstance(value, list) and value else [value]
    content_type = _resolve_multipart_content_type(field_name, schema, encoding)
    files: list[dict[str, Any]] = []

    for offset, item in enumerate(values):
        placeholder_name = field_name if len(values) == 1 else f"{field_name}_{offset + 1}"
        if isinstance(item, str) and item.strip():
            file_path = item.strip()
            source_type = "placeholder" if PLACEHOLDER_PATTERN.fullmatch(file_path) else "path"
        else:
            file_path = f"{{{{{placeholder_name}}}}}"
            source_type = "placeholder"
        file_name = Path(file_path).name if source_type == "path" else placeholder_name
        files.append(
            {
                "field_name": str(field_name),
                "source_type": source_type,
                "file_path": file_path,
                "file_name": file_name or str(field_name),
                "content_type": content_type,
                "base64_content": "",
                "enabled": True,
                "order": order_start + offset,
            }
        )
    return files


def split_multipart_example(
    example: Any,
    schema: dict[str, Any] | None = None,
    encoding: dict[str, Any] | None = None,
) -> tuple[dict[str, Any], list[dict[str, Any]]]:
    text_parts: dict[str, Any] = {}
    files: list[dict[str, Any]] = []
    properties = (schema or {}).get("properties") or {}
    required_names = set((schema or {}).get("required") or [])

    if isinstance(example, dict):
        for key, value in example.items():
            property_schema = properties.get(key) if isinstance(properties, dict) else {}
            if is_binary_schema(property_schema if isinstance(property_schema, dict) else None):
                files.extend(
                    _build_multipart_file_entries(
                        field_name=str(key),
                        value=value,
                        schema=property_schema if isinstance(property_schema, dict) else None,
                        encoding=encoding,
                        order_start=len(files),
                    )
                )
            else:
                text_parts[str(key)] = value
        return text_parts, files

    for index, (name, property_schema) in enumerate((properties or {}).items()):
        if not isinstance(property_schema, dict):
            continue
        if is_binary_schema(property_schema):
            files.extend(
                _build_multipart_file_entries(
                    field_name=str(name),
                    value=None,
                    schema=property_schema,
                    encoding=encoding,
                    order_start=len(files),
                )
            )
        elif name in required_names or property_schema.get("default") not in (None, ""):
            text_parts[str(name)] = property_schema.get("default", "")
    return text_parts, files


class VariableResolver:
    def __init__(self, variables: dict[str, Any] | None = None):
        base_variables = variables.copy() if variables else {}
        base_variables.setdefault("timestamp", int(__import__("time").time()))
        base_variables.setdefault("timestamp_ms", int(__import__("time").time() * 1000))
        base_variables.setdefault("uuid", str(uuid.uuid4()))
        base_variables.setdefault("random_int", random.randint(100000, 999999))
        self.variables = base_variables

    def resolve(self, value: Any) -> Any:
        if isinstance(value, dict):
            return {key: self.resolve(item) for key, item in value.items()}
        if isinstance(value, list):
            return [self.resolve(item) for item in value]
        if isinstance(value, str):
            return self.resolve_string(value)
        return value

    def resolve_string(self, value: str) -> Any:
        matched = PLACEHOLDER_PATTERN.fullmatch(value.strip())
        if matched:
            return self.variables.get(matched.group(1), value)

        def replace(match: re.Match[str]) -> str:
            key = match.group(1)
            replacement = self.variables.get(key, match.group(0))
            if isinstance(replacement, (dict, list)):
                return json.dumps(replacement, ensure_ascii=False)
            return str(replacement)

        return PLACEHOLDER_PATTERN.sub(replace, value)


def build_request_url(base_url: str, request_url: str) -> str:
    if request_url.startswith(("http://", "https://")):
        return request_url
    if not base_url:
        return request_url
    return urljoin(base_url.rstrip("/") + "/", request_url.lstrip("/"))


def extract_json_path(data: Any, path: str) -> Any:
    current = data
    for part in path.split("."):
        if isinstance(current, list):
            if not part.isdigit():
                return None
            index = int(part)
            if index >= len(current):
                return None
            current = current[index]
        elif isinstance(current, dict):
            if part not in current:
                return None
            current = current[part]
        else:
            return None
    return current


def evaluate_assertions(
    assertions: list[dict[str, Any]],
    status_code: int,
    response_text: str,
    response_json: Any,
) -> tuple[list[dict[str, Any]], bool]:
    if not assertions:
        return [], True

    results: list[dict[str, Any]] = []
    all_passed = True

    for index, assertion in enumerate(assertions, start=1):
        assertion_type = assertion.get("type")
        expected = assertion.get("expected")
        operator = assertion.get("operator", "equals")
        path = assertion.get("path")
        passed = False
        actual: Any = None
        message = ""

        if assertion_type == "status_code":
            actual = status_code
            passed = actual == int(expected)
            message = f"状态码应为 {expected}"
        elif assertion_type == "body_contains":
            actual = response_text
            passed = str(expected) in response_text
            message = f"响应体包含 {expected}"
        elif assertion_type == "json_path":
            actual = extract_json_path(response_json, str(path)) if response_json is not None and path else None
            if operator == "contains":
                passed = str(expected) in str(actual)
            elif operator == "not_equals":
                passed = actual != expected
            else:
                passed = actual == expected
            message = f"JSONPath {path} {operator} {expected}"
        elif assertion_type == "header":
            actual = None
            passed = False
            message = f"暂不支持 header 断言: {path or ''}"
        else:
            actual = None
            passed = False
            message = f"未知断言类型: {assertion_type}"

        results.append(
            {
                "index": index,
                "type": assertion_type,
                "operator": operator,
                "path": path,
                "expected": expected,
                "actual": actual,
                "passed": passed,
                "message": message,
            }
        )
        all_passed = all_passed and passed

    return results, all_passed


def import_requests_from_document(file_path: str) -> DocumentImportResult:
    suffix = Path(file_path).suffix.lower()

    structured_result = parse_structured_document(file_path)
    if structured_result is not None:
        return structured_result

    if suffix in TEXT_EXTENSIONS:
        markdown = load_text_document(file_path)
        requests = extract_requests_from_markdown(markdown)
        return DocumentImportResult("text_markdown", requests, False, "已从文本接口文档中抽取接口定义")

    if suffix in NATIVE_DOCUMENT_EXTENSIONS:
        markdown = extract_text_document_content(file_path)
        requests = extract_requests_from_markdown(markdown)
        return DocumentImportResult("native_document", requests, False, "已通过本地文档解析器抽取接口定义")

    if suffix in MARKER_EXTENSIONS:
        markdown = convert_document_with_marker(file_path)
        requests = extract_requests_from_markdown(markdown)
        return DocumentImportResult("marker_markdown", requests, True, "已通过 marker 转换文档并抽取接口定义")

    raise ValueError(f"暂不支持的接口文档格式: {suffix or 'unknown'}")


def load_document_content_for_ai(file_path: str) -> tuple[str, str, bool]:
    suffix = Path(file_path).suffix.lower()

    if suffix in STRUCTURED_EXTENSIONS | TEXT_EXTENSIONS:
        return load_text_document(file_path), "text_document", False

    if suffix in NATIVE_DOCUMENT_EXTENSIONS:
        return extract_text_document_content(file_path), "native_document", False

    if suffix in MARKER_EXTENSIONS:
        return convert_document_with_marker(file_path), "marker_markdown", True

    raise ValueError(f"Unsupported document format for AI parsing: {suffix or 'unknown'}")


def parse_structured_document(file_path: str) -> DocumentImportResult | None:
    suffix = Path(file_path).suffix.lower()
    content = Path(file_path).read_text(encoding="utf-8", errors="ignore")

    if suffix == ".json":
        parsed = json.loads(content)
        if is_openapi_document(parsed):
            return DocumentImportResult("openapi", parse_openapi_document(parsed), False, "已按 OpenAPI/Swagger 结构解析")
        if is_postman_collection(parsed):
            return DocumentImportResult("postman", parse_postman_collection(parsed), False, "已按 Postman Collection 结构解析")
        return None

    if suffix in {".yaml", ".yml"}:
        parsed = yaml.safe_load(content)
        if isinstance(parsed, dict) and is_openapi_document(parsed):
            return DocumentImportResult("openapi", parse_openapi_document(parsed), False, "已按 OpenAPI/Swagger YAML 解析")
        return None

    if looks_like_openapi_text(content):
        parsed = yaml.safe_load(content)
        if isinstance(parsed, dict) and is_openapi_document(parsed):
            return DocumentImportResult("openapi", parse_openapi_document(parsed), False, "已按 OpenAPI 文本解析")

    return None


def is_openapi_document(data: Any) -> bool:
    return isinstance(data, dict) and (("openapi" in data and "paths" in data) or ("swagger" in data and "paths" in data))


def looks_like_openapi_text(content: str) -> bool:
    lowered = content.lower()
    return "paths:" in lowered and ("openapi:" in lowered or "swagger:" in lowered)


def is_postman_collection(data: Any) -> bool:
    return isinstance(data, dict) and "item" in data and "info" in data


def resolve_openapi_ref(spec: dict[str, Any], value: Any) -> Any:
    if not isinstance(value, dict):
        return value
    ref = str(value.get("$ref") or "").strip()
    if not ref.startswith("#/"):
        return value
    current: Any = spec
    for part in ref[2:].split("/"):
        if not isinstance(current, dict) or part not in current:
            return value
        current = current[part]
    return current


def normalize_auth_variable_name(value: str, fallback: str = "token") -> str:
    normalized = re.sub(r"[^a-zA-Z0-9]+", "_", str(value or "")).strip("_").lower()
    if not normalized:
        return fallback
    if normalized in {"bearer", "bearer_auth", "access_token", "authorization", "token"}:
        return "token"
    return normalized


def extract_openapi_security_schemes(spec: dict[str, Any]) -> dict[str, dict[str, Any]]:
    if spec.get("swagger"):
        raw_schemes = spec.get("securityDefinitions") or {}
    else:
        raw_schemes = ((spec.get("components") or {}).get("securitySchemes") or {})

    resolved_schemes: dict[str, dict[str, Any]] = {}
    for scheme_name, scheme_value in raw_schemes.items():
        resolved = resolve_openapi_ref(spec, scheme_value)
        if isinstance(resolved, dict):
            resolved_schemes[str(scheme_name)] = resolved
    return resolved_schemes


def build_openapi_auth_spec(
    scheme_name: str,
    scheme_definition: dict[str, Any] | None,
) -> dict[str, Any] | None:
    if not isinstance(scheme_definition, dict):
        return None

    scheme_type = str(scheme_definition.get("type") or "").lower()
    if scheme_type == "http":
        http_scheme = str(scheme_definition.get("scheme") or "").lower()
        if http_scheme == "basic":
            return {
                "auth": {
                    **default_auth_spec(),
                    "auth_type": "basic",
                    "username": "",
                    "password": "",
                }
            }
        if http_scheme == "bearer":
            return {
                "auth": {
                    **default_auth_spec(),
                    "auth_type": "bearer",
                    "header_name": "Authorization",
                    "bearer_prefix": "Bearer",
                    "token_variable": normalize_auth_variable_name(scheme_name, "token"),
                }
            }
        return None

    if scheme_type == "apikey":
        key_name = str(scheme_definition.get("name") or scheme_name or "")
        key_location = str(scheme_definition.get("in") or "header").lower()
        token_variable = normalize_auth_variable_name(key_name or scheme_name, "token")
        if key_location == "cookie":
            return {
                "auth": {
                    **default_auth_spec(),
                    "auth_type": "cookie",
                    "cookie_name": key_name or "token",
                    "token_variable": token_variable,
                }
            }
        return {
            "auth": {
                **default_auth_spec(),
                "auth_type": "api_key",
                "api_key_name": key_name,
                "api_key_in": key_location if key_location in {"header", "query", "cookie"} else "header",
                "token_variable": token_variable,
            }
        }

    if scheme_type in {"oauth2", "openidconnect"}:
        return {
            "auth": {
                **default_auth_spec(),
                "auth_type": "bearer",
                "header_name": "Authorization",
                "bearer_prefix": "Bearer",
                "token_variable": normalize_auth_variable_name(scheme_name, "token"),
            }
        }

    if scheme_type == "basic":
        return {
            "auth": {
                **default_auth_spec(),
                "auth_type": "basic",
                "username": "",
                "password": "",
            }
        }

    return None


def resolve_openapi_security_auth(
    *,
    spec: dict[str, Any],
    operation: dict[str, Any],
    security_schemes: dict[str, dict[str, Any]],
) -> dict[str, Any]:
    if "security" in operation:
        security_requirements = operation.get("security")
    else:
        security_requirements = spec.get("security")

    if security_requirements == []:
        return {"auth": default_auth_spec()}
    if not isinstance(security_requirements, list):
        return {}

    for requirement in security_requirements:
        if requirement == {}:
            return {"auth": default_auth_spec()}
        if not isinstance(requirement, dict):
            continue
        for scheme_name in requirement.keys():
            auth_payload = build_openapi_auth_spec(str(scheme_name), security_schemes.get(str(scheme_name)))
            if auth_payload:
                return auth_payload
    return {}


def extract_swagger_consumes(spec: dict[str, Any], path_item: dict[str, Any], operation: dict[str, Any]) -> list[str]:
    for owner in (operation, path_item, spec):
        consumes = owner.get("consumes") if isinstance(owner, dict) else None
        if isinstance(consumes, list) and consumes:
            return [str(item).strip() for item in consumes if str(item).strip()]
    return []


def _build_file_spec_from_swagger_parameter(
    parameter: dict[str, Any],
    *,
    content_type: str = "application/octet-stream",
    order: int = 0,
) -> dict[str, Any]:
    field_name = str(parameter.get("name") or f"file_{order}")
    example = extract_parameter_example(parameter)
    file_path = str(example or f"{{{{{field_name}}}}}")
    source_type = "placeholder" if PLACEHOLDER_PATTERN.fullmatch(file_path) else "path"
    file_name = Path(file_path).name if source_type == "path" else field_name
    return {
        "field_name": field_name,
        "source_type": source_type,
        "file_path": file_path,
        "file_name": file_name or field_name,
        "content_type": content_type,
        "base64_content": "",
        "enabled": True,
        "order": order,
    }


def parse_swagger_parameter_body_spec(
    spec: dict[str, Any],
    *,
    parameters: list[dict[str, Any]],
    consumes: list[str] | None = None,
) -> dict[str, Any]:
    consumes = [str(item).lower() for item in (consumes or []) if str(item).strip()]
    body_parameter = next(
        (
            parameter
            for parameter in parameters
            if isinstance(parameter, dict) and str(parameter.get("in") or "").lower() == "body"
        ),
        None,
    )
    if body_parameter:
        schema = resolve_openapi_ref(spec, body_parameter.get("schema") or {})
        example = extract_parameter_example(body_parameter, schema if isinstance(schema, dict) else None)
        content_type = consumes[0] if consumes else "application/json"
        if is_graphql_content_type(content_type):
            if isinstance(example, dict):
                return {
                    "body_mode": "graphql",
                    "graphql_query": str(example.get("query") or ""),
                    "graphql_operation_name": str(example.get("operationName") or ""),
                    "graphql_variables": example.get("variables") if isinstance(example.get("variables"), dict) else {},
                }
            return {"body_mode": "graphql", "graphql_query": str(example or "")}
        if is_xml_content_type(content_type):
            return {
                "body_mode": "xml",
                "xml_text": example if isinstance(example, str) else json.dumps(example, ensure_ascii=False),
            }
        if content_type in TEXT_CONTENT_TYPES or content_type.startswith("text/"):
            return {
                "body_mode": "raw",
                "raw_text": example if isinstance(example, str) else json.dumps(example, ensure_ascii=False),
            }
        if (
            content_type == "application/octet-stream"
            or content_type.startswith("image/")
            or content_type.startswith("audio/")
            or content_type.startswith("video/")
        ):
            return {"body_mode": "binary", "binary_base64": ""}
        return {
            "body_mode": "json",
            "body_json": example if isinstance(example, (dict, list)) else build_example_from_schema(schema if isinstance(schema, dict) else {}),
        }

    form_parameters = [
        parameter
        for parameter in parameters
        if isinstance(parameter, dict) and str(parameter.get("in") or "").lower() == "formdata"
    ]
    if not form_parameters:
        return {}

    multipart_content_type = next((item for item in consumes if item.startswith("multipart/")), "")
    if any(str(parameter.get("type") or "").lower() == "file" for parameter in form_parameters) or multipart_content_type:
        multipart_parts: dict[str, Any] = {}
        files: list[dict[str, Any]] = []
        for parameter in form_parameters:
            field_type = str(parameter.get("type") or "").lower()
            if field_type == "file":
                files.append(
                    _build_file_spec_from_swagger_parameter(
                        parameter,
                        content_type="application/octet-stream",
                        order=len(files),
                    )
                )
                continue
            multipart_parts[str(parameter.get("name") or "")] = extract_parameter_example(parameter)
        return {
            "body_mode": "multipart",
            "multipart_parts": named_items_from_mapping(multipart_parts),
            "files": files,
        }

    form_fields: dict[str, Any] = {}
    for parameter in form_parameters:
        form_fields[str(parameter.get("name") or "")] = extract_parameter_example(parameter)
    return {
        "body_mode": "urlencoded",
        "form_fields": named_items_from_mapping(form_fields),
    }


def parse_openapi_document(spec: dict[str, Any]) -> list[ParsedRequestData]:
    requests: list[ParsedRequestData] = []
    security_schemes = extract_openapi_security_schemes(spec)
    base_url = ""
    servers = spec.get("servers") or []
    if servers and isinstance(servers[0], dict):
        base_url = str(servers[0].get("url") or "").strip()
    elif spec.get("swagger"):
        host = str(spec.get("host") or "").strip()
        base_path = str(spec.get("basePath") or "").strip()
        schemes = spec.get("schemes") or ["http"]
        scheme = str(schemes[0]).strip() if schemes else "http"
        if host:
            normalized_base_path = ""
            if base_path:
                normalized_base_path = base_path if base_path.startswith("/") else f"/{base_path}"
            base_url = f"{scheme}://{host}{normalized_base_path}"

    for path, path_item in (spec.get("paths") or {}).items():
        path_item = resolve_openapi_ref(spec, path_item)
        if not isinstance(path_item, dict):
            continue
        path_level_params = [
            resolve_openapi_ref(spec, parameter)
            for parameter in (path_item.get("parameters") or [])
        ]
        for method, operation in path_item.items():
            upper_method = str(method).upper()
            operation = resolve_openapi_ref(spec, operation)
            if upper_method not in HTTP_METHODS or not isinstance(operation, dict):
                continue

            parameters = list(path_level_params) + [
                resolve_openapi_ref(spec, parameter)
                for parameter in (operation.get("parameters") or [])
            ]
            params: dict[str, Any] = {}
            headers: dict[str, Any] = {}
            cookies: dict[str, Any] = {}
            for parameter in parameters:
                if not isinstance(parameter, dict):
                    continue
                name = parameter.get("name")
                location = parameter.get("in")
                if not name or not location:
                    continue
                example = extract_parameter_example(parameter)
                if location == "query":
                    params[name] = example
                elif location == "header":
                    headers[name] = example
                elif location == "cookie":
                    cookies[name] = example

            body_type = "none"
            body: Any = {}
            request_spec = build_request_spec(
                method=upper_method,
                url=normalize_request_url(base_url, path),
                headers=headers,
                query=params,
                cookies=cookies,
            )
            request_spec.update(
                resolve_openapi_security_auth(
                    spec=spec,
                    operation=operation,
                    security_schemes=security_schemes,
                )
            )
            request_body = resolve_openapi_ref(spec, operation.get("requestBody") or {})
            if isinstance(request_body, dict) and request_body:
                request_spec.update(parse_request_body_spec(request_body))
                body_type, body = request_spec_to_legacy_body(request_spec)
            elif spec.get("swagger"):
                request_spec.update(
                    parse_swagger_parameter_body_spec(
                        spec,
                        parameters=parameters,
                        consumes=extract_swagger_consumes(spec, path_item, operation),
                    )
                )
                body_type, body = request_spec_to_legacy_body(request_spec)

            success_status = extract_success_status(operation.get("responses") or {})
            tags = operation.get("tags") or []
            request_name = operation.get("summary") or operation.get("operationId") or f"{upper_method} {path}"
            final_url = normalize_request_url(base_url, path)

            requests.append(
                ParsedRequestData(
                    name=request_name,
                    method=upper_method,
                    url=final_url,
                    description=str(operation.get("description") or ""),
                    headers=headers,
                    params=params,
                    body_type=body_type,
                    body=body,
                    assertions=[{"type": "status_code", "expected": success_status}],
                    request_spec=request_spec,
                    assertion_specs=[{"assertion_type": "status_code", "expected_number": success_status}],
                    collection_name=tags[0] if tags else guess_collection_name_from_path(path),
                )
            )

    return requests


def parse_postman_collection(collection: dict[str, Any]) -> list[ParsedRequestData]:
    requests: list[ParsedRequestData] = []

    def walk_items(
        items: list[dict[str, Any]],
        folder_name: str | None = None,
        inherited_auth: dict[str, Any] | None = None,
    ):
        for item in items:
            if not isinstance(item, dict):
                continue
            item_auth = _resolve_postman_auth_payload(
                item.get("auth") if "auth" in item else None,
                inherited_auth,
            )
            if "item" in item and isinstance(item["item"], list):
                walk_items(item["item"], item.get("name") or folder_name, item_auth)
                continue

            request_data = item.get("request") or {}
            method = str(request_data.get("method") or "GET").upper()
            if method not in HTTP_METHODS:
                continue

            url = parse_postman_url(request_data.get("url"))
            params = parse_postman_query(request_data.get("url"))
            headers = parse_postman_headers(request_data.get("header") or [])
            request_spec = build_request_spec(
                method=method,
                url=url,
                headers=headers,
                query=params,
            )
            request_spec.update(
                _resolve_postman_auth_payload(
                    request_data.get("auth") if "auth" in request_data else None,
                    item_auth,
                )
            )
            request_spec.update(parse_postman_body_spec(request_data.get("body") or {}))
            body_type, body = request_spec_to_legacy_body(request_spec)
            requests.append(
                ParsedRequestData(
                    name=item.get("name") or f"{method} {url}",
                    method=method,
                    url=url,
                    description=str(item.get("description") or ""),
                    headers=headers,
                    params=params,
                    body_type=body_type,
                    body=body,
                    assertions=[{"type": "status_code", "expected": 200}],
                    request_spec=request_spec,
                    assertion_specs=[{"assertion_type": "status_code", "expected_number": 200}],
                    collection_name=folder_name,
                )
            )

    walk_items(collection.get("item") or [], inherited_auth=parse_postman_auth(collection.get("auth") or {}))
    return requests


def _resolve_postman_auth_payload(
    auth_data: Any,
    inherited_auth: dict[str, Any] | None = None,
) -> dict[str, Any]:
    if auth_data is None:
        return dict(inherited_auth or {})
    parsed_auth = parse_postman_auth(auth_data)
    if parsed_auth:
        return parsed_auth
    return dict(inherited_auth or {})


def parse_postman_url(url_data: Any) -> str:
    if isinstance(url_data, str):
        return url_data
    if isinstance(url_data, dict):
        raw = url_data.get("raw")
        if raw:
            return str(raw)
        host = "".join(url_data.get("host") or [])
        path = "/".join(url_data.get("path") or [])
        protocol = url_data.get("protocol")
        if host:
            base = f"{protocol}://{host}" if protocol else host
            return f"{base}/{path}".rstrip("/")
        return f"/{path}".replace("//", "/")
    return "/"


def parse_postman_headers(header_list: list[dict[str, Any]]) -> dict[str, Any]:
    headers: dict[str, Any] = {}
    for header in header_list:
        if isinstance(header, dict) and header.get("key") and not header.get("disabled"):
            headers[str(header["key"])] = header.get("value", "")
    return headers


def parse_postman_query(url_data: Any) -> dict[str, Any]:
    if isinstance(url_data, dict):
        query_items = url_data.get("query") or []
        if isinstance(query_items, list):
            return {
                str(item.get("key")): item.get("value", "")
                for item in query_items
                if isinstance(item, dict) and item.get("key") and not item.get("disabled")
            }
        raw = str(url_data.get("raw") or "")
        if raw:
            return dict(parse_qsl(raw.split("?", 1)[1])) if "?" in raw else {}
    if isinstance(url_data, str) and "?" in url_data:
        return dict(parse_qsl(url_data.split("?", 1)[1]))
    return {}


def parse_postman_auth(auth_data: dict[str, Any]) -> dict[str, Any]:
    if not isinstance(auth_data, dict):
        return {}
    auth_type = str(auth_data.get("type") or "").lower()
    if auth_type in {"noauth", "none"}:
        return {"auth": default_auth_spec()}
    if auth_type == "bearer":
        token = ""
        for item in auth_data.get("bearer") or []:
            if isinstance(item, dict) and item.get("key") == "token":
                token = str(item.get("value") or "")
                break
        return {"auth": {**default_auth_spec(), "auth_type": "bearer", "token_value": token}}
    if auth_type == "basic":
        username = ""
        password = ""
        for item in auth_data.get("basic") or []:
            if not isinstance(item, dict):
                continue
            if item.get("key") == "username":
                username = str(item.get("value") or "")
            if item.get("key") == "password":
                password = str(item.get("value") or "")
        return {"auth": {**default_auth_spec(), "auth_type": "basic", "username": username, "password": password}}
    if auth_type == "apikey":
        payload = {**default_auth_spec(), "auth_type": "api_key"}
        for item in auth_data.get("apikey") or []:
            if not isinstance(item, dict):
                continue
            key = str(item.get("key") or "")
            value = item.get("value")
            if key == "key":
                payload["api_key_name"] = str(value or "")
            elif key == "value":
                payload["api_key_value"] = str(value or "")
            elif key == "in":
                payload["api_key_in"] = str(value or "header")
        return {"auth": payload}
    return {}


def parse_postman_body_spec(body_data: dict[str, Any]) -> dict[str, Any]:
    mode = body_data.get("mode")
    if mode == "raw":
        raw = body_data.get("raw", "")
        language = str(((body_data.get("options") or {}).get("raw") or {}).get("language") or "").lower()
        if language == "graphql":
            parsed = try_parse_json(raw)
            if isinstance(parsed, dict):
                return {
                    "body_mode": "graphql",
                    "graphql_query": str(parsed.get("query") or ""),
                    "graphql_operation_name": str(parsed.get("operationName") or ""),
                    "graphql_variables": parsed.get("variables") if isinstance(parsed.get("variables"), dict) else {},
                }
            return {"body_mode": "graphql", "graphql_query": str(raw or "")}
        if language == "xml":
            return {"body_mode": "xml", "xml_text": str(raw or "")}
        try:
            return {"body_mode": "json", "body_json": json.loads(raw)}
        except Exception:  # noqa: BLE001
            return {"body_mode": "raw", "raw_text": str(raw or "")}
    if mode == "urlencoded":
        values = {}
        for item in body_data.get("urlencoded") or []:
            if isinstance(item, dict) and item.get("key"):
                values[str(item["key"])] = item.get("value", "")
        return {"body_mode": "urlencoded", "form_fields": named_items_from_mapping(values)}
    if mode == "formdata":
        text_parts: dict[str, Any] = {}
        files: list[dict[str, Any]] = []
        for item in body_data.get("formdata") or []:
            if not isinstance(item, dict) or not item.get("key"):
                continue
            if str(item.get("type") or "text") == "file":
                files.append(
                    {
                        "field_name": str(item["key"]),
                        "source_type": "path",
                        "file_path": str(item.get("src") or ""),
                        "file_name": str(item.get("src") or item["key"]).split("/")[-1].split("\\")[-1],
                        "content_type": str(item.get("contentType") or "application/octet-stream"),
                        "base64_content": "",
                        "enabled": not bool(item.get("disabled")),
                        "order": len(files),
                    }
                )
            else:
                text_parts[str(item["key"])] = item.get("value", "")
        return {
            "body_mode": "multipart",
            "multipart_parts": named_items_from_mapping(text_parts),
            "files": files,
        }
    return {"body_mode": "none"}


def parse_postman_body(body_data: dict[str, Any]) -> tuple[str, Any]:
    request_spec = build_request_spec(method="POST", url="/")
    request_spec.update(parse_postman_body_spec(body_data))
    return request_spec_to_legacy_body(request_spec)


def load_text_document(file_path: str) -> str:
    path = Path(file_path)
    if path.suffix.lower() in {".html", ".htm"}:
        soup = BeautifulSoup(path.read_text(encoding="utf-8", errors="ignore"), "html.parser")
        return soup.get_text("\n")
    return path.read_text(encoding="utf-8", errors="ignore")


def extract_text_document_content(file_path: str) -> str:
    suffix = Path(file_path).suffix.lower()
    if suffix in {".html", ".htm"}:
        return load_text_document(file_path)
    if suffix == ".docx":
        return extract_docx_text(file_path)
    if suffix == ".pptx":
        return extract_pptx_text(file_path)
    if suffix == ".xlsx":
        return extract_xlsx_text(file_path)
    if suffix == ".epub":
        return extract_epub_text(file_path)
    return load_text_document(file_path)


def extract_docx_text(file_path: str) -> str:
    from docx import Document

    document = Document(file_path)
    lines: list[str] = []
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if text:
            lines.append(text)
    for table in document.tables:
        for row in table.rows:
            values = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if values:
                lines.append(" | ".join(values))
    return "\n".join(lines)


def extract_pptx_text(file_path: str) -> str:
    from pptx import Presentation

    presentation = Presentation(file_path)
    lines: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        lines.append(f"# Slide {index}")
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text and text.strip():
                lines.append(text.strip())
    return "\n".join(lines)


def extract_xlsx_text(file_path: str) -> str:
    from openpyxl import load_workbook

    workbook = load_workbook(file_path, data_only=True)
    lines: list[str] = []
    for sheet in workbook.worksheets:
        lines.append(f"# Sheet {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            values = [str(value).strip() for value in row if value is not None and str(value).strip()]
            if values:
                lines.append(" | ".join(values))
    return "\n".join(lines)


def extract_epub_text(file_path: str) -> str:
    from ebooklib import ITEM_DOCUMENT, epub

    book = epub.read_epub(file_path)
    lines: list[str] = []
    for item in book.get_items():
        if item.get_type() != ITEM_DOCUMENT:
            continue
        soup = BeautifulSoup(item.get_body_content(), "html.parser")
        text = soup.get_text("\n").strip()
        if text:
            lines.append(text)
    return "\n\n".join(lines)


def convert_document_with_marker(file_path: str) -> str:
    marker_single = find_marker_single()
    if marker_single:
        return run_marker_cli(marker_single, file_path)

    if importlib.util.find_spec("marker") is not None:
        try:
            from marker.converters.pdf import PdfConverter
            from marker.models import create_model_dict
            from marker.output import text_from_rendered

            converter = PdfConverter(artifact_dict=create_model_dict())
            rendered = converter(file_path)
            text, _, _ = text_from_rendered(rendered)
            return text
        except Exception as exc:  # noqa: BLE001
            raise ValueError(f"marker 文档转换失败: {exc}") from exc

    raise ValueError("当前环境未安装 marker。请先安装 marker-pdf[full]，以支持 PDF、图片等文档导入。")


def find_marker_single() -> str | None:
    configured = os.environ.get("MARKER_SINGLE_PATH")
    if configured and Path(configured).exists():
        return configured

    local_runtime = Path(__file__).resolve().parents[2] / "marker_runtime" / "Scripts" / "marker_single.exe"
    if local_runtime.exists():
        return str(local_runtime)

    local_runtime_cmd = Path(__file__).resolve().parents[2] / "marker_runtime" / "Scripts" / "marker_single"
    if local_runtime_cmd.exists():
        return str(local_runtime_cmd)

    return shutil.which("marker_single")


def run_marker_cli(marker_single: str, file_path: str) -> str:
    with tempfile.TemporaryDirectory(prefix="flytest-marker-") as tmpdir:
        command = [
            marker_single,
            file_path,
            "--output_format",
            "markdown",
            "--output_dir",
            tmpdir,
        ]
        completed = subprocess.run(command, capture_output=True, text=True, check=False)
        if completed.returncode != 0:
            stderr = completed.stderr.strip()
            if "libgobject-2.0-0" in stderr:
                raise ValueError(
                    "marker 在当前 Windows 环境下缺少 WeasyPrint 依赖库，暂时无法直接转换该 Office 文档。"
                    " 当前版本已内置 DOCX/PPTX/XLSX/EPUB/HTML 本地解析兜底，请优先使用这些格式的原生导入；"
                    " PDF 和图片仍可直接通过 marker 导入。"
                )
            raise ValueError(stderr or completed.stdout.strip() or "marker CLI 执行失败")

        stem = Path(file_path).stem
        markdown_path = Path(tmpdir) / f"{stem}.md"
        if not markdown_path.exists():
            nested_markdown_path = Path(tmpdir) / stem / f"{stem}.md"
            if nested_markdown_path.exists():
                markdown_path = nested_markdown_path
            else:
                candidates = list(Path(tmpdir).rglob("*.md"))
                if not candidates:
                    raise ValueError("marker 未生成 Markdown 输出")
                markdown_path = candidates[0]
        return markdown_path.read_text(encoding="utf-8", errors="ignore")


def extract_requests_from_markdown(markdown: str) -> list[ParsedRequestData]:
    requests: list[ParsedRequestData] = []
    seen: set[tuple[str, str]] = set()

    for parsed in extract_requests_from_curl(markdown):
        dedupe_key = (parsed.method, parsed.url)
        if dedupe_key not in seen:
            seen.add(dedupe_key)
            requests.append(parsed)

    for parsed in extract_requests_from_structured_text(markdown):
        dedupe_key = (parsed.method, parsed.url)
        if dedupe_key not in seen:
            seen.add(dedupe_key)
            requests.append(parsed)

    headings = list(HEADING_PATTERN.finditer(markdown))
    code_blocks = list(CODE_BLOCK_PATTERN.finditer(markdown))

    for match in ENDPOINT_PATTERN.finditer(markdown):
        method = match.group("method").upper()
        url = match.group("url").strip()
        dedupe_key = (method, url)
        if dedupe_key in seen:
            continue

        heading = nearest_heading(match.start(), headings)
        body_type = "none"
        body: Any = {}
        next_json_block = find_next_json_block(match.end(), code_blocks)
        if next_json_block and method not in {"GET", "DELETE", "HEAD", "OPTIONS"}:
            parsed = try_parse_json(next_json_block)
            if parsed is not None:
                body_type = "json"
                body = parsed
            else:
                body_type = "raw"
                body = next_json_block.strip()

        request_name = heading or f"{method} {url}"
        requests.append(
            ParsedRequestData(
                name=request_name,
                method=method,
                url=url,
                description=request_name,
                body_type=body_type,
                body=body,
                assertions=[{"type": "status_code", "expected": 200}],
                collection_name=heading,
            )
        )
        seen.add(dedupe_key)

    if not requests:
        raise ValueError("未能从接口文档中识别到接口，请优先上传 Swagger/OpenAPI/Postman 文件，或包含清晰 cURL 示例的接口文档。")

    return requests


def extract_requests_from_structured_text(markdown: str) -> list[ParsedRequestData]:
    requests: list[ParsedRequestData] = []
    candidate = create_structured_candidate()

    for raw_line in markdown.splitlines():
        line = normalize_structured_line(raw_line)
        if not line:
            finalized = finalize_structured_candidate(candidate)
            if finalized:
                requests.append(finalized)
            candidate = create_structured_candidate()
            continue

        label, value = split_structured_label(line)
        field_type = match_structured_field(label) if label else None

        if field_type == "name":
            finalized = finalize_structured_candidate(candidate)
            if finalized:
                requests.append(finalized)
                candidate = create_structured_candidate()
            candidate["name"] = value or candidate["name"]
            continue

        if field_type == "method":
            method = extract_method_from_text(value)
            if method:
                candidate["method"] = method
            continue

        if field_type == "url":
            url = extract_url_from_text(value)
            if url:
                candidate["url"] = url
            continue

        if field_type == "headers":
            candidate["headers"].update(parse_key_value_payload(value))
            continue

        if field_type == "query":
            candidate["query_params"].update(parse_parameter_payload(value))
            continue

        if field_type == "request_params":
            candidate["request_params"].update(parse_parameter_payload(value))
            continue

        if field_type == "body":
            body_type, body = parse_body_payload(value)
            if body_type != "none":
                candidate["body_type"] = body_type
                candidate["body"] = body
            continue

        if field_type == "success":
            candidate["success_notes"].append(value)
            continue

        if field_type == "description":
            candidate["description_lines"].append(value)
            continue

        inline_method = extract_method_from_text(line)
        inline_url = extract_url_from_text(line)
        if inline_method and inline_url:
            finalized = finalize_structured_candidate(candidate)
            if finalized:
                requests.append(finalized)
                candidate = create_structured_candidate()
            candidate["method"] = inline_method
            candidate["url"] = inline_url
            if not candidate["name"]:
                candidate["name"] = line
            continue

        if not candidate["name"] and looks_like_structured_name(line):
            candidate["name"] = line
            continue

        candidate["description_lines"].append(line)

    finalized = finalize_structured_candidate(candidate)
    if finalized:
        requests.append(finalized)

    deduped_requests: list[ParsedRequestData] = []
    seen: set[tuple[str, str]] = set()
    for item in requests:
        key = (item.method, item.url)
        if key in seen:
            continue
        seen.add(key)
        deduped_requests.append(item)
    return deduped_requests


def create_structured_candidate() -> dict[str, Any]:
    return {
        "name": "",
        "method": "",
        "url": "",
        "headers": {},
        "query_params": {},
        "request_params": {},
        "body_type": "none",
        "body": {},
        "description_lines": [],
        "success_notes": [],
    }


def normalize_structured_line(line: str) -> str:
    normalized = line.strip()
    normalized = re.sub(r"^#{1,6}\s*", "", normalized)
    normalized = re.sub(r"^[\-\*\u2022]\s*", "", normalized)
    normalized = re.sub(r"^\d+\.\s*", "", normalized)
    normalized = re.sub("^\\d+\u3001\\s*", "", normalized)
    return normalized.strip()


def split_structured_label(line: str) -> tuple[str | None, str]:
    for separator in ("\uFF1A", ":"):
        if separator in line:
            label, value = line.split(separator, 1)
            return label.strip(), value.strip()
    return None, line


def match_structured_field(label: str | None) -> str | None:
    if not label:
        return None
    normalized = re.sub(r"\s+", "", label).lower()
    field_aliases = {
        "name": {
            "\u63a5\u53e3\u540d\u79f0",
            "\u63a5\u53e3\u540d",
            "api\u540d\u79f0",
            "apiname",
            "\u540d\u79f0",
            "name",
        },
        "method": {
            "\u8bf7\u6c42\u65b9\u5f0f",
            "\u8bf7\u6c42\u65b9\u6cd5",
            "method",
            "httpmethod",
        },
        "url": {
            "\u8bf7\u6c42\u5730\u5740",
            "\u63a5\u53e3\u5730\u5740",
            "\u8bf7\u6c42\u8def\u5f84",
            "\u63a5\u53e3\u8def\u5f84",
            "\u8bf7\u6c42url",
            "\u63a5\u53e3url",
            "url",
            "path",
            "endpoint",
        },
        "headers": {"\u8bf7\u6c42\u5934", "header", "headers"},
        "query": {
            "\u67e5\u8be2\u53c2\u6570",
            "query",
            "queryparams",
            "querystring",
            "url\u53c2\u6570",
        },
        "request_params": {
            "\u8bf7\u6c42\u53c2\u6570",
            "\u5165\u53c2",
            "\u53c2\u6570",
            "requestparams",
            "path\u53c2\u6570",
            "pathparams",
        },
        "body": {
            "\u8bf7\u6c42\u4f53",
            "body",
            "body\u53c2\u6570",
            "\u8bf7\u6c42\u62a5\u6587",
            "bodypayload",
        },
        "success": {
            "\u6210\u529f\u54cd\u5e94",
            "\u54cd\u5e94\u7ed3\u679c",
            "\u8fd4\u56de\u7ed3\u679c",
            "\u8fd4\u56de\u793a\u4f8b",
            "\u54cd\u5e94\u793a\u4f8b",
            "\u54cd\u5e94\u5185\u5bb9",
            "\u8fd4\u56de\u5185\u5bb9",
            "\u8fd4\u56de\u5b57\u6bb5",
        },
        "description": {
            "\u63a5\u53e3\u63cf\u8ff0",
            "\u63cf\u8ff0",
            "\u8bf4\u660e",
            "\u7528\u9014",
            "\u529f\u80fd\u8bf4\u660e",
        },
    }
    for field_type, aliases in field_aliases.items():
        if normalized in aliases:
            return field_type
    return None


def looks_like_structured_name(line: str) -> bool:
    if extract_url_from_text(line):
        return False
    if extract_method_from_text(line):
        return False
    if len(line) > 80:
        return False
    return bool(re.search(r"[\u4e00-\u9fffA-Za-z]", line))


def extract_method_from_text(value: str) -> str | None:
    matched = re.search(r"\b(GET|POST|PUT|PATCH|DELETE|HEAD|OPTIONS)\b", value, re.I)
    if not matched:
        return None
    method = matched.group(1).upper()
    return method if method in HTTP_METHODS else None


def extract_url_from_text(value: str) -> str | None:
    matched = URL_IN_TEXT_PATTERN.search(value)
    if not matched:
        return None
    return matched.group(1).rstrip(".,;，；)")


def parse_key_value_payload(value: str) -> dict[str, Any]:
    parsed_json = try_parse_json(value)
    if isinstance(parsed_json, dict):
        return parsed_json

    pairs: dict[str, Any] = {}
    for key, raw_value in re.findall(r"([A-Za-z_][A-Za-z0-9_.-]*)\s*(?:=|:|：)\s*([^,\n，；;]+)", value):
        pairs[key] = raw_value.strip()
    return pairs


def parse_parameter_payload(value: str) -> dict[str, Any]:
    parsed_json = try_parse_json(value)
    if isinstance(parsed_json, dict):
        return parsed_json

    pairs = parse_key_value_payload(value)
    if pairs:
        return pairs

    params: dict[str, Any] = {}
    for token in re.split(r"[,\n，；;、]+", value):
        token = token.strip()
        if not token:
            continue
        matched = re.match(r"([A-Za-z_][A-Za-z0-9_.-]*)", token)
        if not matched:
            continue
        key = matched.group(1)
        params[key] = f"{{{{{key}}}}}"
    return params


def parse_body_payload(value: str) -> tuple[str, Any]:
    parsed_json = try_parse_json(value)
    if parsed_json is not None:
        return "json", parsed_json

    parsed_params = parse_parameter_payload(value)
    if parsed_params:
        return "json", parsed_params

    stripped = value.strip()
    if stripped:
        return "raw", stripped
    return "none", {}


def extract_success_status_from_text(text: str, method: str) -> int:
    matched = STATUS_CODE_TEXT_PATTERN.search(text)
    if matched:
        try:
            return int(matched.group(1))
        except ValueError:
            pass
    if method == "POST" and any(
        keyword in text for keyword in ("\u521b\u5efa", "\u65b0\u589e", "create")
    ):
        return 201
    return 200


def extract_success_assertions(text: str, method: str) -> list[dict[str, Any]]:
    assertions: list[dict[str, Any]] = [{"type": "status_code", "expected": extract_success_status_from_text(text, method)}]
    seen_paths: set[str] = set()
    for path in JSON_PATH_TEXT_PATTERN.findall(text):
        if path in seen_paths:
            continue
        seen_paths.add(path)
        assertions.append(
            {
                "type": "json_path",
                "path": path,
                "operator": "not_equals",
                "expected": None,
            }
        )
        if len(seen_paths) >= 3:
            break
    return assertions


def finalize_structured_candidate(candidate: dict[str, Any]) -> ParsedRequestData | None:
    method = extract_method_from_text(candidate["method"]) or "GET"
    url = extract_url_from_text(candidate["url"])
    if not url:
        return None

    params = dict(candidate["query_params"])
    body_type = candidate["body_type"]
    body = candidate["body"]

    if candidate["request_params"]:
        if method in {"GET", "DELETE", "HEAD", "OPTIONS"}:
            params.update(candidate["request_params"])
        elif body_type == "none":
            body_type = "json"
            body = candidate["request_params"]

    success_text = " ".join(candidate["success_notes"]).strip()
    description_parts = [part for part in candidate["description_lines"] if part]
    if success_text:
        description_parts.append(success_text)
    description = " ".join(description_parts)[:5000]

    name = candidate["name"].strip() if candidate["name"] else ""
    if not name:
        name = f"{method} {url}"

    return ParsedRequestData(
        name=name[:120],
        method=method,
        url=url,
        description=description,
        headers=candidate["headers"],
        params=params,
        body_type=body_type if body_type in {"none", "json", "form", "raw"} else "none",
        body=body if body_type != "none" else {},
        assertions=extract_success_assertions(success_text, method),
        collection_name=guess_collection_name_from_path(url),
    )


def extract_requests_from_curl(markdown: str) -> list[ParsedRequestData]:
    requests: list[ParsedRequestData] = []
    headings = list(HEADING_PATTERN.finditer(markdown))

    for block_match in CODE_BLOCK_PATTERN.finditer(markdown):
        content = block_match.group("content") or ""
        if "curl " not in content:
            continue
        normalized = content.replace("\\\n", " ").replace("\\\r\n", " ").strip()
        try:
            args = shlex.split(normalized)
        except ValueError:
            continue
        if not args or args[0] != "curl":
            continue

        method = "GET"
        url = ""
        headers: dict[str, Any] = {}
        params: dict[str, Any] = {}
        cookies: dict[str, Any] = {}
        body_mode = "none"
        body_json: Any = {}
        raw_text = ""
        xml_text = ""
        graphql_query = ""
        graphql_variables: dict[str, Any] = {}
        form_fields: dict[str, Any] = {}
        multipart_parts: dict[str, Any] = {}
        files: list[dict[str, Any]] = []
        auth = default_auth_spec()
        transport = default_transport_spec()
        data_tokens: list[str] = []
        explicit_method = False
        force_get_query = False
        i = 1
        while i < len(args):
            token = args[i]
            if token in {"-X", "--request"} and i + 1 < len(args):
                method = args[i + 1].upper()
                explicit_method = True
                i += 2
                continue
            if token in {"-I", "--head"}:
                method = "HEAD"
                explicit_method = True
                i += 1
                continue
            if token in {"-H", "--header"} and i + 1 < len(args):
                header_value = args[i + 1]
                if ":" in header_value:
                    key, value = header_value.split(":", 1)
                    headers[key.strip()] = value.strip()
                i += 2
                continue
            if token in {"-A", "--user-agent"} and i + 1 < len(args):
                headers["User-Agent"] = str(args[i + 1])
                i += 2
                continue
            if token in {"-b", "--cookie"} and i + 1 < len(args):
                for cookie_entry in str(args[i + 1]).split(";"):
                    if "=" not in cookie_entry:
                        continue
                    key, value = cookie_entry.split("=", 1)
                    cookies[key.strip()] = value.strip()
                i += 2
                continue
            if token in {"-u", "--user"} and i + 1 < len(args):
                credentials = str(args[i + 1])
                username, _, password = credentials.partition(":")
                auth = {
                    **default_auth_spec(),
                    "auth_type": "basic",
                    "username": username,
                    "password": password,
                }
                i += 2
                continue
            if token in {"-k", "--insecure"}:
                transport["verify_ssl"] = False
                i += 1
                continue
            if token in {"-L", "--location"}:
                transport["follow_redirects"] = True
                i += 1
                continue
            if token in {"-G", "--get"}:
                force_get_query = True
                i += 1
                continue
            if token == "--max-time" and i + 1 < len(args):
                try:
                    transport["timeout_ms"] = max(1, int(float(args[i + 1]) * 1000))
                except (TypeError, ValueError):
                    pass
                i += 2
                continue
            if token in {"-x", "--proxy"} and i + 1 < len(args):
                transport["proxy_url"] = str(args[i + 1])
                i += 2
                continue
            if token == "--cert" and i + 1 < len(args):
                transport["client_cert"] = str(args[i + 1])
                i += 2
                continue
            if token == "--key" and i + 1 < len(args):
                transport["client_key"] = str(args[i + 1])
                i += 2
                continue
            if token in {"-d", "--data", "--data-raw", "--data-binary", "--data-urlencode"} and i + 1 < len(args):
                data_tokens.append(str(args[i + 1]))
                i += 2
                continue
            if token in {"-F", "--form"} and i + 1 < len(args):
                body_mode = "multipart"
                form_value = str(args[i + 1])
                parsed_form_value = parse_curl_form_value(form_value, order=len(files))
                if parsed_form_value:
                    if parsed_form_value.get("kind") == "file":
                        files.append(parsed_form_value["file"])
                    else:
                        multipart_parts[str(parsed_form_value["name"])] = parsed_form_value.get("value", "")
                i += 2
                continue
            if token.startswith(("http://", "https://", "/")):
                url = token
            i += 1

        if not url:
            continue
        if "?" in url:
            base_url, query_string = url.split("?", 1)
            url = base_url
            params = dict(parse_qsl(query_string, keep_blank_values=True))

        if force_get_query and data_tokens:
            for data_token in data_tokens:
                params.update(parse_curl_key_value_pairs(data_token))
        elif data_tokens:
            raw_body = "&".join(data_tokens)
            parsed = try_parse_json(raw_body)
            content_type_header = str(headers.get("Content-Type") or headers.get("content-type") or "")
            form_pairs = parse_curl_data_pairs(data_tokens)
            if is_graphql_content_type(content_type_header):
                body_mode = "graphql"
                if isinstance(parsed, dict):
                    graphql_query = str(parsed.get("query") or "")
                    graphql_variables = parsed.get("variables") if isinstance(parsed.get("variables"), dict) else {}
                else:
                    graphql_query = raw_body
            elif is_xml_content_type(content_type_header):
                body_mode = "xml"
                xml_text = raw_body
            elif parsed is not None:
                body_mode = "json"
                body_json = parsed
            elif str(content_type_header).lower() == "application/x-www-form-urlencoded" or form_pairs is not None:
                body_mode = "urlencoded"
                form_fields = form_pairs or {}
            else:
                body_mode = "raw"
                raw_text = raw_body
            if not explicit_method and method == "GET":
                method = "POST"

        if body_mode == "multipart" and not explicit_method and method == "GET":
            method = "POST"

        request_spec = build_request_spec(
            method=method,
            url=url,
            headers=headers,
            query=params,
            cookies=cookies,
            body_mode=body_mode,
            body_json=body_json,
            raw_text=raw_text,
            xml_text=xml_text,
            graphql_query=graphql_query,
            graphql_variables=graphql_variables,
            form_fields=form_fields,
            multipart_parts=multipart_parts,
            files=files,
            auth=auth,
            transport=transport,
            timeout_ms=int(transport.get("timeout_ms") or DEFAULT_REQUEST_TIMEOUT_MS),
        )
        body_type, body = request_spec_to_legacy_body(request_spec)

        heading = nearest_heading(block_match.start(), headings)
        requests.append(
            ParsedRequestData(
                name=heading or f"{method} {url}",
                method=method,
                url=url,
                headers=headers,
                params=params,
                body_type=body_type,
                body=body,
                assertions=[{"type": "status_code", "expected": 200}],
                request_spec=request_spec,
                assertion_specs=[{"assertion_type": "status_code", "expected_number": 200}],
                collection_name=heading,
            )
        )

    return requests


def parse_curl_key_value_pairs(payload: str) -> dict[str, Any]:
    parsed_pairs = parse_qsl(str(payload or "").lstrip("?"), keep_blank_values=True)
    return {key: value for key, value in parsed_pairs if key}


def parse_curl_data_pairs(payloads: list[str]) -> dict[str, Any] | None:
    merged: dict[str, Any] = {}
    found_pair = False
    for payload in payloads:
        pairs = parse_curl_key_value_pairs(payload)
        if not pairs:
            return None
        merged.update(pairs)
        found_pair = True
    return merged if found_pair else None


def parse_curl_form_value(form_value: str, *, order: int = 0) -> dict[str, Any] | None:
    if "=" not in form_value:
        return None
    key, raw_value = form_value.split("=", 1)
    key = key.strip()
    raw_value = raw_value.strip()
    if not key:
        return None
    if raw_value.startswith("@"):
        file_value = raw_value[1:]
        parts = file_value.split(";")
        file_path = parts[0].strip()
        options: dict[str, str] = {}
        for option in parts[1:]:
            if "=" not in option:
                continue
            option_name, option_value = option.split("=", 1)
            options[option_name.strip().lower()] = option_value.strip()
        source_type = "placeholder" if PLACEHOLDER_PATTERN.fullmatch(file_path) else "path"
        file_name = options.get("filename") or (Path(file_path).name if source_type == "path" else key)
        return {
            "kind": "file",
            "file": {
                "field_name": key,
                "source_type": source_type,
                "file_path": file_path or f"{{{{{key}}}}}",
                "file_name": file_name or key,
                "content_type": options.get("type") or "application/octet-stream",
                "base64_content": "",
                "enabled": True,
                "order": order,
            },
        }
    return {"kind": "field", "name": key, "value": raw_value}


def nearest_heading(position: int, headings: list[re.Match[str]]) -> str | None:
    current = None
    for heading in headings:
        if heading.start() >= position:
            break
        current = heading.group("title").strip()
    return current


def find_next_json_block(position: int, code_blocks: list[re.Match[str]]) -> str | None:
    for block in code_blocks:
        if block.start() < position:
            continue
        language = (block.group("lang") or "").lower()
        content = block.group("content") or ""
        if language in {"json", "javascript", "js"}:
            return content
        if try_parse_json(content) is not None:
            return content
        return None
    return None


def try_parse_json(text: str) -> Any | None:
    stripped = text.strip()
    if not stripped:
        return None
    try:
        return json.loads(stripped)
    except Exception:  # noqa: BLE001
        return None


def extract_parameter_example(parameter: dict[str, Any], schema_override: dict[str, Any] | None = None) -> Any:
    if "x-example" in parameter:
        return parameter.get("x-example")
    if "example" in parameter:
        return parameter.get("example")
    if str(parameter.get("type") or "").lower() == "file":
        return f"{{{{{parameter.get('name') or 'file'}}}}}"
    schema = schema_override if isinstance(schema_override, dict) else parameter.get("schema") or {}
    if not schema and parameter.get("type"):
        schema = {"type": parameter.get("type")}
        if "default" in parameter:
            schema["default"] = parameter.get("default")
        if "enum" in parameter:
            schema["enum"] = parameter.get("enum")
        if "items" in parameter:
            schema["items"] = parameter.get("items")
    if "x-example" in schema:
        return schema.get("x-example")
    if "example" in schema:
        return schema.get("example")
    if "default" in schema:
        return schema.get("default")
    enum_values = schema.get("enum") or []
    if enum_values:
        return enum_values[0]
    schema_example = build_example_from_schema(schema)
    return schema_example


def parse_request_body_spec(request_body: dict[str, Any]) -> dict[str, Any]:
    content = request_body.get("content") or {}
    for content_type, content_item in content.items():
        lowered = str(content_type).lower()
        example = extract_content_example(content_item) if isinstance(content_item, dict) else {}
        if is_graphql_content_type(lowered) or (is_json_content_type(lowered) and looks_like_graphql_payload(example)):
            if isinstance(example, dict):
                return {
                    "body_mode": "graphql",
                    "graphql_query": str(example.get("query") or ""),
                    "graphql_operation_name": str(example.get("operationName") or ""),
                    "graphql_variables": example.get("variables") if isinstance(example.get("variables"), dict) else {},
                }
            return {
                "body_mode": "graphql",
                "graphql_query": str(example or ""),
                "graphql_operation_name": "",
                "graphql_variables": {},
            }
        if is_json_content_type(lowered):
            return {
                "body_mode": "json",
                "body_json": example,
            }
        if lowered == "application/x-www-form-urlencoded":
            return {
                "body_mode": "urlencoded",
                "form_fields": named_items_from_mapping(example if isinstance(example, dict) else {}),
            }
        if lowered.startswith("multipart/"):
            schema = content_item.get("schema") if isinstance(content_item, dict) else {}
            encoding = content_item.get("encoding") if isinstance(content_item, dict) else {}
            multipart_parts, files = split_multipart_example(
                example,
                schema if isinstance(schema, dict) else None,
                encoding if isinstance(encoding, dict) else None,
            )
            return {
                "body_mode": "multipart",
                "multipart_parts": named_items_from_mapping(multipart_parts),
                "files": files,
            }
        if is_xml_content_type(lowered):
            return {
                "body_mode": "xml",
                "xml_text": example if isinstance(example, str) else json.dumps(example, ensure_ascii=False),
            }
        if lowered in TEXT_CONTENT_TYPES or lowered.startswith("text/"):
            return {
                "body_mode": "raw",
                "raw_text": example if isinstance(example, str) else json.dumps(example, ensure_ascii=False),
            }
        if (
            lowered == "application/octet-stream"
            or lowered.startswith("image/")
            or lowered.startswith("audio/")
            or lowered.startswith("video/")
        ):
            return {
                "body_mode": "binary",
                "binary_base64": "",
            }
    return {"body_mode": "none"}


def parse_request_body(request_body: dict[str, Any]) -> tuple[str, Any]:
    request_spec = build_request_spec(method="POST", url="/")
    request_spec.update(parse_request_body_spec(request_body))
    return request_spec_to_legacy_body(request_spec)


def extract_content_example(content_item: dict[str, Any]) -> Any:
    if "example" in content_item:
        return content_item["example"]
    examples = content_item.get("examples") or {}
    if isinstance(examples, dict) and examples:
        first_example = next(iter(examples.values()))
        if isinstance(first_example, dict) and "value" in first_example:
            return first_example["value"]
        return first_example
    schema = content_item.get("schema") or {}
    return build_example_from_schema(schema)


def looks_like_graphql_payload(example: Any) -> bool:
    if not isinstance(example, dict):
        return False
    query_text = example.get("query")
    return isinstance(query_text, str) and bool(query_text.strip())


def build_example_from_schema(schema: dict[str, Any]) -> Any:
    if not isinstance(schema, dict):
        return {}
    if "example" in schema:
        return schema["example"]
    schema_type = schema.get("type")
    if schema_type == "object":
        return {key: build_example_from_schema(value) for key, value in (schema.get("properties") or {}).items()}
    if schema_type == "array":
        return [build_example_from_schema(schema.get("items") or {})]
    if schema_type == "integer":
        return schema.get("default", 0)
    if schema_type == "number":
        return schema.get("default", 0)
    if schema_type == "boolean":
        return schema.get("default", False)
    return schema.get("default", "")


def extract_success_status(responses: dict[str, Any]) -> int:
    for key in responses.keys():
        if str(key).startswith("2"):
            try:
                return int(str(key))
            except ValueError:
                continue
    return 200


def normalize_request_url(base_url: str, path: str) -> str:
    if not base_url:
        return path if path.startswith("/") else f"/{path}"
    return base_url.rstrip("/") + "/" + path.lstrip("/")


def guess_collection_name_from_path(path: str) -> str | None:
    segments = [segment for segment in path.split("/") if segment]
    return segments[0] if segments else None
